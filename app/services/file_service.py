import os
import uuid
import csv
import io
from flask import current_app
from werkzeug.utils import secure_filename

def allowed_file(filename):
    """Check if the file extension is allowed"""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in current_app.config['ALLOWED_EXTENSIONS']

def save_file(file, chapter_id):
    """Save the file to the upload folder and return the path"""
    # Generate unique filename to prevent collisions
    original_filename = secure_filename(file.filename)
    file_extension = original_filename.rsplit('.', 1)[1].lower() if '.' in original_filename else ''
    unique_filename = f"{uuid.uuid4().hex}.{file_extension}"
    
    # Create directory structure based on chapter ID
    upload_dir = os.path.join(current_app.config['UPLOAD_FOLDER'], f'chapter_{chapter_id}')
    os.makedirs(upload_dir, exist_ok=True)
    
    file_path = os.path.join(upload_dir, unique_filename)
    relative_path = os.path.join(f'chapter_{chapter_id}', unique_filename)
    
    # Save the file
    file.save(file_path)
    
    return relative_path

def get_file_path(relative_path):
    """Get the absolute file path from the relative path stored in the database"""
    return os.path.join(current_app.config['UPLOAD_FOLDER'], relative_path)

def extract_text_from_file(file_path):
    """Extract text from various file types"""
    file_extension = file_path.rsplit('.', 1)[1].lower() if '.' in file_path else ''
    
    if file_extension == 'pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension in ['doc', 'docx']:
        return extract_text_from_doc(file_path)
    elif file_extension in ['ppt', 'pptx']:
        return extract_text_from_ppt(file_path)
    elif file_extension in ['xls', 'xlsx']:
        return extract_text_from_excel(file_path)
    elif file_extension == 'csv':
        return extract_text_from_csv(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")

def extract_text_from_pdf(file_path):
    """Extract text from PDF files"""
    try:
        # Using PyPDF2 (you'll need to add this to requirements.txt)
        from PyPDF2 import PdfReader
        
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        
        return text
    except Exception as e:
        current_app.logger.error(f"Error extracting text from PDF: {str(e)}")
        return f"Error extracting text: {str(e)}"

def extract_text_from_doc(file_path):
    """Extract text from DOC/DOCX files"""
    try:
        # Using python-docx for DOCX
        if file_path.endswith('.docx'):
            import docx
            
            doc = docx.Document(file_path)
            full_text = []
            
            # Extract text from paragraphs
            for para in doc.paragraphs:
                full_text.append(para.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text)
                    full_text.append(" | ".join(row_text))
            
            return "\n".join(full_text)
        else:
            # For DOC files (older format), use docx2txt if available
            try:
                import docx2txt
                text = docx2txt.process(file_path)
                return text
            except ImportError:
                # Fallback method using basic binary read
                try:
                    import re
                    with open(file_path, 'rb') as file:
                        content = file.read()
                        # Simple extraction of readable text (limited capability)
                        text_pattern = re.compile(b'[\x20-\x7E]{4,}')  # Printable ASCII sequences
                        matches = text_pattern.findall(content)
                        extracted_text = [match.decode('ascii', errors='ignore') for match in matches]
                        return "\n".join(extracted_text)
                except:
                    current_app.logger.warning("Unable to extract text from DOC file")
                    return "Text extraction from DOC files requires docx2txt library"
    except Exception as e:
        current_app.logger.error(f"Error extracting text from DOC/DOCX: {str(e)}")
        return f"Error extracting text: {str(e)}"

def extract_text_from_ppt(file_path):
    """Extract text from PPT/PPTX files"""
    try:
        # Using python-pptx for PPTX
        if file_path.endswith('.pptx'):
            from pptx import Presentation
            
            text = []
            prs = Presentation(file_path)
            for slide in prs.slides:
                slide_text = []
                slide_text.append(f"--- Slide {len(text) + 1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                text.append("\n".join(slide_text))
            
            return "\n\n".join(text)
        else:
            # For PPT files, use basic binary extraction
            try:
                import re
                with open(file_path, 'rb') as file:
                    content = file.read()
                    # Simple extraction of readable text (limited capability)
                    text_pattern = re.compile(b'[\x20-\x7E]{4,}')  # Printable ASCII sequences
                    matches = text_pattern.findall(content)
                    extracted_text = [match.decode('ascii', errors='ignore') for match in matches]
                    
                    # Filter out common binary file junk
                    filtered_text = []
                    for line in extracted_text:
                        if not any(marker in line.lower() for marker in ['microsoft', 'powerpoint', 'font', 'times', 'arial']):
                            filtered_text.append(line)
                    
                    return "\n".join(filtered_text)
            except Exception as e:
                current_app.logger.warning(f"Failed to extract text from PPT: {e}")
                return "Full text extraction from PPT files is limited without additional libraries"
    except Exception as e:
        current_app.logger.error(f"Error extracting text from PPT/PPTX: {str(e)}")
        return f"Error extracting text: {str(e)}"

def extract_text_from_excel(file_path):
    """Extract text from Excel (XLS/XLSX) files"""
    try:
        import pandas as pd
        
        # Read all sheets
        excel_file = pd.ExcelFile(file_path)
        sheet_names = excel_file.sheet_names
        
        text_output = []
        
        # Process each sheet
        for sheet_name in sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            # Add sheet name as header
            text_output.append(f"Sheet: {sheet_name}")
            
            # Check if dataframe is empty
            if df.empty:
                text_output.append("(Empty sheet)")
                continue
            
            # Convert dataframe to string representation
            # First, convert all values to strings to handle mixed types
            df_string = df.astype(str)
            
            # Get headers
            headers = " | ".join(df_string.columns)
            text_output.append(headers)
            
            # Add separator
            text_output.append("-" * len(headers))
            
            # Process rows
            for _, row in df_string.iterrows():
                text_output.append(" | ".join(row.values))
            
            # Add spacing between sheets
            text_output.append("\n")
        
        return "\n".join(text_output)
    except Exception as e:
        current_app.logger.error(f"Error extracting text from Excel file: {str(e)}")
        return f"Error extracting text: {str(e)}"

def extract_text_from_csv(file_path):
    """Extract text from CSV files"""
    try:
        text_output = []
        
        # Try to detect encoding
        encodings = ['utf-8', 'latin-1', 'iso-8859-1', 'cp1252']
        
        # Try different encodings
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as csv_file:
                    # Try to sniff the dialect
                    try:
                        dialect = csv.Sniffer().sniff(csv_file.read(4096))
                        csv_file.seek(0)
                    except:
                        dialect = csv.excel  # Fallback to standard dialect
                        csv_file.seek(0)
                    
                    reader = csv.reader(csv_file, dialect)
                    
                    # Read header
                    try:
                        headers = next(reader)
                        text_output.append(" | ".join(headers))
                        text_output.append("-" * len(" | ".join(headers)))
                    except StopIteration:
                        text_output.append("(Empty CSV file)")
                        break
                    
                    # Read rows
                    for row in reader:
                        text_output.append(" | ".join(row))
                
                # If we got here without errors, break the encoding loop
                break
            except UnicodeDecodeError:
                if encoding == encodings[-1]:  # Last encoding attempt
                    text_output.append(f"Unable to decode CSV file with available encodings")
                continue
            except Exception as e:
                text_output.append(f"Error reading CSV: {str(e)}")
                break
        
        return "\n".join(text_output)
    except Exception as e:
        current_app.logger.error(f"Error extracting text from CSV file: {str(e)}")
        return f"Error extracting text: {str(e)}"

